"""
PowerPoint Exporter Module - Corrected Version
Generates investor pitch deck presentations
"""

import os
import logging
from typing import Dict, Any, Optional, List
from datetime import datetime

logger = logging.getLogger(__name__)

# Try to import python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
PPTX_AVAILABLE = True

# Color scheme
COLORS = {
    'primary': RGBColor(0x1a, 0x36, 0x5d),  # Dark blue
    'secondary': RGBColor(0x38, 0xa1, 0x69),  # Green
    'accent': RGBColor(0x66, 0x5d, 0xc3),  # Purple
    'text': RGBColor(0x2d, 0x37, 0x48),  # Dark gray
    'light': RGBColor(0xf7, 0xfa, 0xfc),  # Light gray
    'white': RGBColor(0xff, 0xff, 0xff),
}



def generate_pptx_report(
    output_path: str,
    company_name: str,
    industry: str,
    valuation_data: Dict[str, Any],
    assumptions: Dict[str, Any],
    historical_data: Optional[Dict[str, Any]] = None,
    commentary: Optional[Dict[str, str]] = None
) -> bool:
    """
    Generate a PowerPoint presentation
    
    Args:
        output_path: Path to save the PPTX
        company_name: Name of the company
        industry: Industry classification
        valuation_data: Valuation metrics
        assumptions: Model assumptions
        historical_data: Historical financials (optional)
        commentary: AI-generated commentary (optional)
    
    Returns:
        True if successful, False otherwise
    """
    if not PPTX_AVAILABLE:
        logger.error("python-pptx not available for PowerPoint generation")
        return False
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title
        _add_title_slide(prs, company_name, industry)
        
        # Slide 2: Executive Summary
        _add_summary_slide(prs, company_name, valuation_data, commentary)
        
        # Slide 3: Valuation Overview
        _add_valuation_slide(prs, valuation_data)
        
        # Slide 4: Key Assumptions
        _add_assumptions_slide(prs, assumptions)
        
        # Slide 5: Investment Thesis (if commentary available)
        if commentary:
            _add_thesis_slide(prs, commentary)
        
        # Slide 6: Disclaimer
        _add_disclaimer_slide(prs)
        
        # Save presentation
        prs.save(output_path)
        
        logger.info(f"PowerPoint presentation generated: {output_path}")
        return True
        
    except Exception as e:
        logger.error(f"Error generating PowerPoint: {e}")
        return False


def _add_title_slide(prs: 'Presentation', company_name: str, industry: str):
    """Add title slide"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary']
    background.line.fill.background()
    
    # Company name
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = company_name
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        Inches(12.333), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = f"Financial Model & Valuation | {industry.title()}"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLORS['light']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5),
        Inches(12.333), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.text = datetime.now().strftime("%B %Y")
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = COLORS['light']
    date_para.alignment = PP_ALIGN.CENTER


def _add_summary_slide(prs: 'Presentation', company_name: str, valuation_data: Dict, commentary: Optional[Dict]):
    """Add executive summary slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    _add_slide_title(slide, "Executive Summary")
    
    # Key metrics boxes
    metrics = [
        ("Enterprise Value", f"₹{valuation_data.get('enterprise_value', 0):,.0f} Cr", COLORS['primary']),
        ("Equity Value", f"₹{valuation_data.get('equity_value', 0):,.0f} Cr", COLORS['secondary']),
        ("Share Price", f"₹{valuation_data.get('share_price', 0):,.2f}", COLORS['accent']),
        ("WACC", f"{valuation_data.get('wacc', 0) * 100:.1f}%", COLORS['primary']),
    ]
    
    box_width = Inches(2.8)
    box_height = Inches(1.5)
    start_x = Inches(0.8)
    y = Inches(1.8)
    gap = Inches(0.3)
    
    for i, (label, value, color) in enumerate(metrics):
        x = start_x + (i * (box_width + gap))
        
        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Label
        label_box = slide.shapes.add_textbox(x, y + Inches(0.2), box_width, Inches(0.4))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = COLORS['white']
        label_para.alignment = PP_ALIGN.CENTER
        
        # Value
        value_box = slide.shapes.add_textbox(x, y + Inches(0.6), box_width, Inches(0.7))
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = value
        value_para.font.size = Pt(24)
        value_para.font.bold = True
        value_para.font.color.rgb = COLORS['white']
        value_para.alignment = PP_ALIGN.CENTER
    
    # Commentary section
    if commentary and commentary.get('investment_thesis'):
        thesis_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.8),
            Inches(11.733), Inches(2.5)
        )
        thesis_frame = thesis_box.text_frame
        thesis_frame.word_wrap = True
        
        # Thesis header
        p1 = thesis_frame.paragraphs[0]
        p1.text = "Investment Thesis"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLORS['primary']
        
        # Thesis content
        p2 = thesis_frame.add_paragraph()
        p2.text = commentary.get('investment_thesis', '')
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORS['text']
        p2.space_before = Pt(12)


def _add_valuation_slide(prs: 'Presentation', valuation_data: Dict):
    """Add valuation details slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    _add_slide_title(slide, "DCF Valuation")
    
    # Valuation waterfall description
    items = [
        ("PV of Forecast Cash Flows", valuation_data.get('pv_fcf', 0)),
        ("PV of Terminal Value", valuation_data.get('pv_terminal', 0)),
        ("Enterprise Value", valuation_data.get('enterprise_value', 0)),
        ("Less: Net Debt", valuation_data.get('net_debt', 0)),
        ("Equity Value", valuation_data.get('equity_value', 0)),
    ]
    
    y = Inches(2)
    for label, value in items:
        # Label
        label_box = slide.shapes.add_textbox(Inches(1), y, Inches(5), Inches(0.5))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(16)
        label_para.font.color.rgb = COLORS['text']
        
        # Value
        value_box = slide.shapes.add_textbox(Inches(7), y, Inches(4), Inches(0.5))
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = f"₹{value:,.0f} Cr"
        value_para.font.size = Pt(16)
        value_para.font.bold = True
        value_para.font.color.rgb = COLORS['primary']
        value_para.alignment = PP_ALIGN.RIGHT
        
        y += Inches(0.7)


def _add_assumptions_slide(prs: 'Presentation', assumptions: Dict):
    """Add key assumptions slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    _add_slide_title(slide, "Key Assumptions")
    
    # Assumptions in two columns
    left_items = [
        ("Revenue Growth", f"{assumptions.get('revenue_growth', 0) * 100:.1f}%"),
        ("EBITDA Margin", f"{assumptions.get('ebitda_margin', 0) * 100:.1f}%"),
        ("Tax Rate", f"{assumptions.get('tax_rate', 0) * 100:.1f}%"),
        ("CapEx (% of Revenue)", f"{assumptions.get('capex_ratio', 0) * 100:.1f}%"),
    ]
    
    right_items = [
        ("Terminal Growth", f"{assumptions.get('terminal_growth', 0) * 100:.1f}%"),
        ("Risk-Free Rate", f"{assumptions.get('risk_free_rate', 0.07) * 100:.1f}%"),
        ("Equity Risk Premium", f"{assumptions.get('equity_risk_premium', 0.06) * 100:.1f}%"),
        ("Beta", f"{assumptions.get('beta', 1.0):.2f}"),
    ]
    
    # Left column
    y = Inches(2)
    for label, value in left_items:
        _add_assumption_row(slide, Inches(1), y, label, value)
        y += Inches(0.8)
    
    # Right column
    y = Inches(2)
    for label, value in right_items:
        _add_assumption_row(slide, Inches(7), y, label, value)
        y += Inches(0.8)


def _add_assumption_row(slide, x, y, label, value):
    """Helper to add assumption row"""
    # Label
    label_box = slide.shapes.add_textbox(x, y, Inches(3.5), Inches(0.5))
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = label
    label_para.font.size = Pt(14)
    label_para.font.color.rgb = COLORS['text']
    
    # Value
    value_box = slide.shapes.add_textbox(x + Inches(3.5), y, Inches(1.5), Inches(0.5))
    value_frame = value_box.text_frame
    value_para = value_frame.paragraphs[0]
    value_para.text = value
    value_para.font.size = Pt(14)
    value_para.font.bold = True
    value_para.font.color.rgb = COLORS['secondary']
    value_para.alignment = PP_ALIGN.RIGHT


def _add_thesis_slide(prs: 'Presentation', commentary: Dict):
    """Add investment thesis slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    _add_slide_title(slide, "Investment Analysis")
    
    y = Inches(1.8)
    
    sections = [
        ("Investment Thesis", commentary.get('investment_thesis', '')),
        ("Key Risks", commentary.get('key_risks', '')),
        ("Recommendation", commentary.get('recommendation', '')),
    ]
    
    for title, content in sections:
        if not content:
            continue
        
        # Section title
        title_box = slide.shapes.add_textbox(Inches(1), y, Inches(11), Inches(0.4))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(14)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['primary']
        y += Inches(0.5)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(1), y, Inches(11), Inches(1))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_para = content_frame.paragraphs[0]
        content_para.text = content
        content_para.font.size = Pt(12)
        content_para.font.color.rgb = COLORS['text']
        y += Inches(1.2)


def _add_disclaimer_slide(prs: 'Presentation'):
    """Add disclaimer slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    _add_slide_title(slide, "Disclaimer")
    
    disclaimer_text = """This presentation has been prepared for informational purposes only and does not constitute an offer to sell or a solicitation of an offer to buy any securities.

The projections and valuations contained herein are based on assumptions and historical data that may not reflect future performance. Past performance is not indicative of future results.

Investors should conduct their own independent analysis and consult with professional advisors before making any investment decisions.

Generated by AI Financial Modeler"""
    
    text_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(11.333), Inches(4)
    )
    text_frame = text_box.textframe
    text_frame.word_wrap = True
    para = text_frame.paragraphs[0]
    para.text = disclaimer_text
    para.font.size = Pt(12)
    para.font.color.rgb = COLORS['text']
    para.line_spacing = 1.5


def _add_slide_title(slide, title_text: str):
    """Add title to slide"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(12.333), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']


def create_presentation(company_name: str, excel_path: str, output_path: str) -> bool:
    """
    Wrapper function to create PowerPoint presentation from Excel model
    
    Args:
        company_name: Name of the company
        excel_path: Path to the generated Excel model
        output_path: Path to save the PowerPoint presentation
    
    Returns:
        True if successful, False otherwise
    """
    if not PPTX_AVAILABLE:
        logger.error("python-pptx not available for PowerPoint generation")
        return False
    
    try:
        import openpyxl
        
        # Load the Excel workbook
        wb = openpyxl.load_workbook(excel_path)
        
        # Extract industry from Summary sheet (Row 8, Column 3)
        industry = "General"
        if "Summary" in wb.sheetnames:
            summary = wb["Summary"]
            industry_cell = summary.cell(8, 3).value
            if industry_cell:
                industry = str(industry_cell)
        
        # Extract valuation data from Valuation sheet
        valuation_data = {
            'enterprise_value': 0,
            'equity_value': 0,
            'share_price': 0,
            'wacc': 0.10,
            'pv_fcf': 0,
            'pv_terminal': 0,
            'net_debt': 0,
        }
        
        if "Valuation" in wb.sheetnames:
            val_sheet = wb["Valuation"]
            try:
                # Read specific cells
                wacc_val = val_sheet.cell(19, 3).value
                if wacc_val and isinstance(wacc_val, (int, float)):
                    valuation_data['wacc'] = float(wacc_val) if wacc_val <= 1 else wacc_val / 100
                
                ev_val = val_sheet.cell(35, 3).value
                if ev_val and isinstance(ev_val, (int, float)):
                    valuation_data['enterprise_value'] = float(ev_val)
                
                equity_val = val_sheet.cell(37, 3).value
                if equity_val and isinstance(equity_val, (int, float)):
                    valuation_data['equity_value'] = float(equity_val)
                
                price_val = val_sheet.cell(40, 3).value
                if price_val and isinstance(price_val, (int, float)):
                    valuation_data['share_price'] = float(price_val)
                
                # Try to get PV values if available (common locations)
                pv_fcf = val_sheet.cell(33, 3).value
                if pv_fcf and isinstance(pv_fcf, (int, float)):
                    valuation_data['pv_fcf'] = float(pv_fcf)
                
                pv_term = val_sheet.cell(34, 3).value
                if pv_term and isinstance(pv_term, (int, float)):
                    valuation_data['pv_terminal'] = float(pv_term)
            except Exception as e:
                logger.warning(f"Error reading valuation data: {e}")
        
        # Extract assumptions from Assumptions sheet
        assumptions = {
            'revenue_growth': 0.10,
            'ebitda_margin': 0.20,
            'tax_rate': 0.25,
            'terminal_growth': 0.03,
            'risk_free_rate': 0.07,
            'equity_risk_premium': 0.06,
            'beta': 1.0,
            'capex_ratio': 0.04,
        }
        
        if "Assumptions" in wb.sheetnames:
            assump_sheet = wb["Assumptions"]
            try:
                keyword_map = {
                    'ebitda margin': 'ebitda_margin',
                    'terminal growth': 'terminal_growth',
                    'risk-free rate': 'risk_free_rate',
                    'equity risk premium': 'equity_risk_premium',
                    'beta': 'beta',
                    'tax rate': 'tax_rate',
                    'revenue growth': 'revenue_growth',
                    'capex % of revenue': 'capex_ratio',
                }
                
                for row in range(1, 41):
                    label = assump_sheet.cell(row, 2).value
                    if label and isinstance(label, str):
                        label_lower = label.lower()
                        value_cell = assump_sheet.cell(row, 3).value
                        
                        if value_cell and isinstance(value_cell, (int, float)):
                            val = float(value_cell)
                            # Convert percentages to decimals
                            if val > 1 and any(x in label_lower for x in ["rate", "growth", "margin", "premium", "capex"]):
                                val = val / 100
                            
                            for keyword, key in keyword_map.items():
                                if keyword in label_lower:
                                    assumptions[key] = val
                                    break
            except Exception as e:
                logger.warning(f"Error reading assumptions: {e}")
        
        wb.close()
        
        # Generate the PowerPoint
        return generate_pptx_report(
            output_path=output_path,
            company_name=company_name,
            industry=industry,
            valuation_data=valuation_data,
            assumptions=assumptions,
            historical_data=None,
            commentary=None
        )
        
    except Exception as e:
        logger.error(f"Error creating PowerPoint from Excel: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return False


def is_available() -> bool:
    """Check if PowerPoint export is available"""
    return PPTX_AVAILABLE
